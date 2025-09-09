# rag_pipeline.py (diagnostic + improved loaders + optional OCR)
import os, re, io, hashlib
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, Tuple

# ---- loaders ----
import fitz  # PyMuPDF for PDFs
import docx  # python-docx
from pptx import Presentation  # python-pptx

# Optional OCR deps (loaded lazily only if enable_ocr=True)
try:
    import pytesseract  # requires local Tesseract install
    from PIL import Image
except Exception:  # don't hard fail if not installed
    pytesseract = None
    Image = None

# ---- vector store ----
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions

# ---- http for Ollama / optional OpenAI ----
import requests


# ---------------------------
# Text utils
# ---------------------------
def normalize_ws(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200, min_chars: int = 40) -> List[str]:
    """
    Simple sliding window chunker over the whole text with overlap.
    Drops tiny chunks (< min_chars) to avoid noise.
    """
    text = normalize_ws(text)
    if not text:
        return []
    chunks, start, n = [], 0, len(text)
    while start < n:
        end = min(n, start + chunk_size)
        seg = text[start:end]
        if len(seg) >= min_chars:
            chunks.append(seg)
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks


# ---------------------------
# File loaders (richer)
# ---------------------------
def _pptx_collect_shape_text(shape) -> List[str]:
    """Recursively collect text from pptx shapes (handles groups/placeholders)."""
    parts: List[str] = []
    # Group shapes
    if hasattr(shape, "shapes"):
        for shp in shape.shapes:
            parts.extend(_pptx_collect_shape_text(shp))
    # Text frames
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        tf = shape.text_frame
        # text_frame.text gives concatenated text; also gather paragraphs/runs for robustness
        if tf.text:
            parts.append(tf.text)
        else:
            for p in tf.paragraphs:
                runs = [r.text for r in p.runs if r.text]
                if runs:
                    parts.append("".join(runs))
    # Alt text (sometimes used)
    alt = getattr(shape, "alternative_text", "") or ""
    if alt.strip():
        parts.append(alt.strip())
    return parts

def load_pptx(path: str) -> List[Tuple[str, Dict[str, Any]]]:
    out, pres = [], Presentation(path)
    for i, slide in enumerate(pres.slides):
        parts: List[str] = []
        # title + placeholders
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(slide.shapes.title.text)
        # all shapes (including groups)
        for shp in slide.shapes:
            parts.extend(_pptx_collect_shape_text(shp))
        # notes (optional; uncomment to include)
        # if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
        #     nt = slide.notes_slide.notes_text_frame.text
        #     if nt:
        #         parts.append(nt)
        text = normalize_ws("\n".join([p for p in parts if p]))
        if text:
            out.append((text, {"type": "pptx", "slide": i + 1}))
    return out

def load_docx(path: str) -> List[Tuple[str, Dict[str, Any]]]:
    d = docx.Document(path)
    parts: List[str] = []
    # paragraphs
    parts.extend([p.text for p in d.paragraphs if p.text])
    # tables
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    text = normalize_ws("\n".join(parts))
    return [(text, {"type": "docx"})] if text else []

def _pdf_page_to_image_bytes(page: fitz.Page, zoom: float = 2.0) -> bytes:
    """Rasterize a page to PNG bytes for OCR."""
    m = fitz.Matrix(zoom, zoom)  # 2x zoom for better OCR
    pix = page.get_pixmap(matrix=m, alpha=False)
    return pix.tobytes("png")

def load_pdf(path: str, enable_ocr: bool = False, ocr_lang: str = "eng", ocr_max_pages: Optional[int] = None) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Extract text from PDF; if a page has no text and enable_ocr=True, run OCR on a rasterized page image.
    Set ocr_max_pages to limit how many pages to OCR (for speed).
    """
    out: List[Tuple[str, Dict[str, Any]]] = []
    doc = fitz.open(path)
    ocr_used = False
    for i, page in enumerate(doc):
        raw = page.get_text("text") or ""
        text = normalize_ws(raw)
        if not text and enable_ocr:
            if pytesseract is None or Image is None:
                # OCR requested but not available
                pass
            else:
                # optional cap on OCR pages
                if ocr_max_pages is None or sum(1 for _, m in out if m.get("ocr")) < ocr_max_pages:
                    try:
                        png = _pdf_page_to_image_bytes(page, zoom=2.0)
                        img = Image.open(io.BytesIO(png))
                        text = normalize_ws(pytesseract.image_to_string(img, lang=ocr_lang) or "")
                        if text:
                            ocr_used = True
                            out.append((text, {"type": "pdf", "page": i + 1, "ocr": True}))
                            continue
                    except Exception:
                        # ignore OCR failures; fall back to empty page
                        pass
        if text:
            out.append((text, {"type": "pdf", "page": i + 1}))
    doc.close()
    return out

def load_file(path: str, enable_ocr: bool = False, ocr_lang: str = "eng", ocr_max_pages: Optional[int] = None) -> List[Tuple[str, Dict[str, Any]]]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":  return load_pdf(path, enable_ocr=enable_ocr, ocr_lang=ocr_lang, ocr_max_pages=ocr_max_pages)
    if ext == ".docx": return load_docx(path)
    if ext == ".pptx": return load_pptx(path)
    return []


# ---------------------------
# Embeddings
# ---------------------------
class OllamaEmbeddings:
    def __init__(self, model: str = "nomic-embed-text:latest", base_url: str = "http://localhost:11434", timeout: int = 120):
        self.model = model
        self.base_url = base_url.rstrip("/")
        self.timeout = timeout

    def embed(self, texts: List[str]) -> List[List[float]]:
        url = f"{self.base_url}/api/embeddings"
        vecs: List[List[float]] = []
        for t in texts:
            try:
                resp = requests.post(url, json={"model": self.model, "prompt": t}, timeout=self.timeout)
                resp.raise_for_status()
            except requests.HTTPError as e:
                if getattr(e, "response", None) is not None and e.response.status_code == 404:
                    raise RuntimeError(
                        "Ollama embeddings endpoint (/api/embeddings) not found. "
                        "Upgrade Ollama and pull an embedding model, e.g.: "
                        "`ollama pull nomic-embed-text:latest`"
                    ) from e
                raise
            data = resp.json()
            v = data.get("embedding")
            if not isinstance(v, list):
                raise RuntimeError("Invalid embedding response from Ollama.")
            vecs.append(v)
        return vecs


class OpenAIEmbeddings:
    def __init__(self, model: str = "text-embedding-3-small", dimensions: Optional[int] = None):
        from openai import OpenAI
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.model = model
        self.dimensions = dimensions

    def embed(self, texts: List[str]) -> List[List[float]]:
        kwargs = {}
        if self.dimensions:
            kwargs["dimensions"] = self.dimensions
        resp = self.client.embeddings.create(model=self.model, input=texts, **kwargs)
        return [d.embedding for d in resp.data]


# ---------------------------
# Vector DB (Chroma)
# ---------------------------
@dataclass
class ChromaConfig:
    path: str = "chroma_db"
    collection: str = "nps_cyber_ta"

class ChromaStore:
    def __init__(self, cfg: ChromaConfig, ef=None):
        os.makedirs(cfg.path, exist_ok=True)
        self.client = chromadb.PersistentClient(path=cfg.path, settings=Settings(allow_reset=False))
        self.collection = self.client.get_or_create_collection(name=cfg.collection, embedding_function=ef)

    def add_with_embeddings(self, ids, docs, metadatas, embeddings):
        self.collection.add(ids=ids, documents=docs, metadatas=metadatas, embeddings=embeddings)

    def add_docs(self, ids, docs, metadatas):
        self.collection.add(ids=ids, documents=docs, metadatas=metadatas)

    def query(self, query_texts=None, query_embeddings=None, n_results=4, where=None):
        kwargs = {"n_results": n_results}
        if where: kwargs["where"] = where
        if query_embeddings is not None:
            kwargs["query_embeddings"] = query_embeddings
        elif query_texts is not None:
            kwargs["query_texts"] = query_texts
        else:
            raise ValueError("Provide query_texts or query_embeddings")
        return self.collection.query(**kwargs)


# ---------------------------
# RAG Pipeline
# ---------------------------
def _slug(s: str) -> str:
    return re.sub(r"[^a-zA-Z0-9]+", "_", s).strip("_").lower()

@dataclass
class RAGConfig:
    repo_dir: str
    chroma_path: str = "chroma_db"
    collection: str = "nps_cyber_ta"          # base name (we append a suffix automatically)
    embed_provider: str = "ollama"            # "ollama" | "openai" | "onnx"
    embed_model: str = "nomic-embed-text:latest"
    embed_dimensions: Optional[int] = None    # only for openai
    ollama_url: str = "http://localhost:11434"
    chunk_size: int = 1000
    chunk_overlap: int = 200
    min_chars: int = 40                       # drop tiny chunks
    # OCR options (PDF)
    enable_ocr: bool = False                  # set True to OCR empty PDF pages
    ocr_lang: str = "eng"
    ocr_max_pages: Optional[int] = None       # limit OCR pages per file (None = no limit)

class RAGPipeline:
    def __init__(self, cfg: RAGConfig):
        self.cfg = cfg

        ef = None
        # decide embedding + canonical dimension (for suffix)
        if cfg.embed_provider == "openai":
            self.embedder = OpenAIEmbeddings(model=cfg.embed_model, dimensions=cfg.embed_dimensions)
            dim = cfg.embed_dimensions or (1536 if "text-embedding-3-small" in cfg.embed_model else 3072)
        elif cfg.embed_provider == "onnx":
            ef = embedding_functions.ONNXMiniLM_L6_V2()  # 384-d
            self.embedder = None
            dim = 384
        else:
            self.embedder = OllamaEmbeddings(model=cfg.embed_model, base_url=cfg.ollama_url)
            dim = None  # unknown/varies; we isolate by provider+model

        suffix = f"{_slug(cfg.embed_provider)}-{_slug(cfg.embed_model)}-{dim or 'na'}"
        coll_name = f"{cfg.collection}-{suffix}"

        self.store = ChromaStore(ChromaConfig(path=cfg.chroma_path, collection=coll_name), ef=ef)

        print(f"[init] RAGPipeline ready")
        print(f"       Repo dir        : {os.path.abspath(cfg.repo_dir)}")
        print(f"       Chroma path     : {os.path.abspath(cfg.chroma_path)}")
        print(f"       Collection name : {coll_name}")
        print(f"       Provider/Model  : {cfg.embed_provider} / {cfg.embed_model}")
        if cfg.embed_provider == "openai":
            print(f"       Dimensions      : {cfg.embed_dimensions or 'default (per model)'}")
        if cfg.embed_provider == "onnx":
            print(f"       ONNX dims       : 384")
        if cfg.embed_provider == "ollama":
            print(f"       Ollama URL      : {cfg.ollama_url}")
        print(f"       Chunking        : size={cfg.chunk_size}, overlap={cfg.chunk_overlap}, min_chars={cfg.min_chars}")
        if cfg.enable_ocr:
            print(f"       OCR             : enabled (lang={cfg.ocr_lang}, max_pages={cfg.ocr_max_pages})")

    def _gen_id(self, path: str, chunk_index: int, salt: str) -> str:
        h = hashlib.sha1()
        h.update(path.encode("utf-8"))
        h.update(str(chunk_index).encode("utf-8"))
        h.update(salt.encode("utf-8"))
        return h.hexdigest()

    def index_repo(self) -> Dict[str, Any]:
        repo = os.path.abspath(self.cfg.repo_dir)
        if not os.path.isdir(repo):
            raise FileNotFoundError(f"Repo directory not found: {repo}")

        supported = (".pdf", ".docx", ".pptx")
        files: List[str] = []
        for root, _, filenames in os.walk(repo):
            for fn in filenames:
                if os.path.splitext(fn)[1].lower() in supported:
                    files.append(os.path.join(root, fn))
        files.sort()
        print(f"[index] Found {len(files)} supported files.")

        total_added = 0
        total_pages = 0
        total_chars = 0
        pdf_empty_pages = 0
        per_file_stats = []

        for path in files:
            ext = os.path.splitext(path)[1].lower()
            stat = os.stat(path)
            salt = str(int(stat.st_mtime))

            try:
                # load raw records from file
                records = load_file(
                    path,
                    enable_ocr=self.cfg.enable_ocr,
                    ocr_lang=self.cfg.ocr_lang,
                    ocr_max_pages=self.cfg.ocr_max_pages
                )
            except Exception as e:
                print(f"[index] ERROR reading {os.path.basename(path)}: {e}")
                continue

            # Stats: pages/slides/units
            pages_units = len(records)
            total_pages += pages_units

            # Chunk per record
            chunks_all: List[str] = []
            empty_units = 0
            for text, meta in records:
                if not text.strip():
                    empty_units += 1
                    continue
                total_chars += len(text)
                chunks = chunk_text(text, self.cfg.chunk_size, self.cfg.chunk_overlap, self.cfg.min_chars)
                chunks_all.extend(chunks)

            # Track empty PDF pages estimate
            if ext == ".pdf":
                pdf_empty_pages += empty_units

            # Prepare metadatas/ids for chunks
            num_chunks = len(chunks_all)
            if num_chunks:
                ids = [self._gen_id(path, i, salt) for i in range(num_chunks)]
                metas = []
                # we keep a simple running chunk_index here (global to file)
                for i in range(num_chunks):
                    md = {
                        "source": os.path.relpath(path, repo).replace("\\", "/"),
                        "chunk_index": i,
                    }
                    metas.append(md)

                # Store
                if self.store.collection._embedding_function is not None:
                    # ONNX via Chroma (no external calls)
                    self.store.add_docs(ids=ids, docs=chunks_all, metadatas=metas)
                else:
                    # External embedder (Ollama or OpenAI)
                    embeds = self.embedder.embed(chunks_all)
                    self.store.add_with_embeddings(ids=ids, docs=chunks_all, metadatas=metas, embeddings=embeds)

                total_added += num_chunks

            # Log per-file
            print(f"[index] {os.path.basename(path)}: units={pages_units}, empty_units={empty_units}, chunks={num_chunks}")

            per_file_stats.append({
                "file": os.path.basename(path),
                "units": pages_units,
                "empty_units": empty_units,
                "chunks": num_chunks
            })

        print("\n[index] SUMMARY")
        print(f"  Files indexed     : {len(files)}")
        print(f"  Units (pages/slides/records): {total_pages}")
        print(f"  Empty PDF units   : {pdf_empty_pages}")
        print(f"  Total text chars  : {total_chars:,}")
        print(f"  Chunks added      : {total_added}")

        # Also print what Chroma thinks the count is now
        try:
            count_now = self.store.collection.count()
            print(f"  Chroma count      : {count_now}")
        except Exception as e:
            print(f"  Chroma count      : (error: {e})")

        return {
            "indexed_files": len(files),
            "indexed_units": total_pages,
            "indexed_chunks": total_added,
            "per_file": per_file_stats
        }

    def retrieve(self, query: str, k: int = 4) -> List[Dict[str, Any]]:
        # If the collection has no embedding_function, embed the query ourselves
        has_ef = getattr(self.store.collection, "_embedding_function", None) is not None
        if not has_ef and getattr(self, "embedder", None) is not None:
            q_vec = self.embedder.embed([query])
            res = self.store.query(query_embeddings=q_vec, n_results=k)
        else:
            res = self.store.query(query_texts=[query], n_results=k)

        docs = res.get("documents", [[]])[0]
        metas = res.get("metadatas", [[]])[0]
        ids = res.get("ids", [[]])[0]
        out = []
        for doc, md, id_ in zip(docs, metas, ids):
            src = md.get("source", "unknown")
            loc = md.get("page") or md.get("slide") or md.get("chunk_index")
            out.append({"id": id_, "text": doc, "source": src, "loc": loc})
        return out

    @staticmethod
    def format_context(snippets: List[Dict[str, Any]]) -> str:
        parts = []
        for i, s in enumerate(snippets, 1):
            loc_str = f"{s['source']}#{s['loc']}"
            txt = s["text"]
            if len(txt) > 800:
                txt = txt[:800] + " ..."
            parts.append(f"[{i}] ({loc_str})\n{txt}")
        return "Use the following retrieved context to answer. If relevant, cite with [#] or (source#loc).\n\n" + "\n\n".join(parts)
